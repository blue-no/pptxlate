#_cli.py
import os

import click
import pptx
from pptx.presentation import Presentation
from tqdm.autonotebook import tqdm

from pptxlate._trans import translate


@click.group()
def cli() -> None:
    pass


@cli.command()
@click.argument("src")
@click.argument("dst")
@click.option("--kana", "-k", is_flag=True, help="Translate kanas.")
@click.option("--num", "-n", is_flag=True, help="Translate numbers.")
@click.option("--alph", "-a", is_flag=True, help="Translate alphabets.")
@click.option("--skip-title", "-st", is_flag=True, help="Skip title texts.")
def han(
    src: str,
    dst: str,
    kana: bool,
    num: bool,
    alph: bool,
    skip_title: bool,
) -> None:
    if True not in (kana, num, alph):
        exit("Set at least one option: --kana, --num, --alph")
    if os.path.exists(dst) and not _confirm_overwrite():
        exit("Aborted.")
    prs: Presentation = pptx.Presentation(src)
    for slide in tqdm(prs.slides):
        translate(
            slide=slide,
            size_old="full",
            size_new="half",
            kana=kana,
            num=num,
            alph=alph,
            skip_title=skip_title,
        )
    try:
        prs.save(dst)
    except PermissionError:
        exit("Failed to save. Please close Power Point files and try again.")
    exit("Done!")


@cli.command()
@click.argument("src")
@click.argument("dst")
@click.option("--kana", "-k", is_flag=True, help="Translate kanas.")
@click.option("--num", "-n", is_flag=True, help="Translate numbers.")
@click.option("--alph", "-a", is_flag=True, help="Translate alphabets.")
@click.option("--skip-title", "-st", is_flag=True, help="Skip title texts.")
def zen(
    src: str,
    dst: str,
    kana: bool,
    num: bool,
    alph: bool,
    skip_title: bool,
) -> None:
    if True not in (kana, num, alph):
        exit("Set at least one option: --kana, --num, --alph")
    if os.path.exists(dst) and not _confirm_overwrite():
        exit("Aborted.")
    prs: Presentation = pptx.Presentation(src)
    for slide in tqdm(prs.slides):
        translate(
            slide=slide,
            size_old="half",
            size_new="full",
            kana=kana,
            num=num,
            alph=alph,
            skip_title=skip_title,
        )
    try:
        prs.save(dst)
    except PermissionError:
        exit("Failed to save. Please close the PowerPoint file and try again.")
    exit("Done!")


def _confirm_overwrite() -> bool:
    ans = None
    while ans is None:
        try:
            ans = click.confirm(
                "File already exists. Do you want to overwrite it?",
                default=None,
            )
        except NameError:
            ans = None
    return ans
