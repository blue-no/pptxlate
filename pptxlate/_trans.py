#_trans.py
from __future__ import annotations

from typing import Literal

from pptx.slide import Slide

from pptxlate._maps import _alph, _kana, _num


def translate(
    slide: Slide,
    size_old: Literal["full", "half"],
    size_new: Literal["full", "half"],
    kana: bool,
    num: bool,
    alph: bool,
    skip_title: bool,
) -> None:
    for shape in slide.shapes:
        if skip_title and shape == slide.shapes.title:
            continue

        paragraphs = []

        if shape.has_text_frame:
            paragraphs += shape.text_frame.paragraphs

        if shape.has_table:
            for cell in shape.table.iter_cells():
                paragraphs += cell.text_frame.paragraphs

        for paragraph in paragraphs:
            for run in paragraph.runs:
                text = run.text
                if kana:
                    text = _replace(text, _kana[size_old], _kana[size_new])
                if num:
                    text = _replace(text, _num[size_old], _num[size_new])
                if alph:
                    text = _replace(text, _alph[size_old], _alph[size_new])
                run.text = text


def _replace(text: str, olds: list[str], news: list[str]) -> str:
    for old, new in zip(olds, news):
        text = text.replace(old, new)
    return text
